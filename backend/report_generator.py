"""
Gerador de relatórios em PDF e Excel
"""
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib import colors
from reportlab.lib.units import cm
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image, KeepTogether
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from datetime import datetime
import pandas as pd
from typing import List, Dict, Any, Optional
import os
import io
import matplotlib
matplotlib.use('Agg')  # Backend sem interface gráfica
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import font_manager
import numpy as np
from openpyxl.drawing.image import Image as ExcelImage
from openpyxl.utils import get_column_letter

class ReportGenerator:
    """Gerador de relatórios"""
    
    def __init__(self, db=None, client_id=None):
        self.styles = getSampleStyleSheet()
        self.db = db
        self.client_id = client_id
        # Carrega cores primeiro (antes de configurar estilos)
        self._load_client_colors()
        self._setup_custom_styles()
    
    def _load_client_colors(self):
        """Carrega cores do cliente do banco de dados"""
        import json
        from .models import Client
        
        # Cores padrão (Converplast)
        self.cores = {
            'primary': '#1a237e',
            'primaryDark': '#0d47a1',
            'primaryLight': '#3949ab',
            'secondary': '#556B2F',
            'secondaryDark': '#4a5d23',
            'secondaryLight': '#6B8E23',
            'masculino': '#1a237e',
            'feminino': '#556B2F',
            'preto': '#000000',
            'dourado': '#FFD700'
        }
        
        # Paleta padrão
        self.paleta = ['#1a237e', '#556B2F', '#3949ab', '#6B8E23', '#0d47a1', '#808000', '#5c6bc0', '#4a5d23', '#9E9E9E', '#757575']
        
        # Se tem client_id e db, busca cores personalizadas
        if self.client_id and self.db:
            try:
                cliente = self.db.query(Client).filter(Client.id == self.client_id).first()
                if cliente and cliente.cores_personalizadas:
                    cores_json = json.loads(cliente.cores_personalizadas)
                    if isinstance(cores_json, dict):
                        self.cores.update(cores_json)
                        # Atualiza paleta baseada nas cores
                        self.paleta = [
                            self.cores.get('primary', '#1a237e'),
                            self.cores.get('secondary', '#556B2F'),
                            self.cores.get('primaryLight', '#3949ab'),
                            self.cores.get('secondaryLight', '#6B8E23'),
                            self.cores.get('primaryDark', '#0d47a1'),
                            self.cores.get('secondaryDark', '#808000'),
                            self.cores.get('primaryLighter', '#5c6bc0'),
                            self.cores.get('secondaryLighter', '#4a5d23'),
                            self.cores.get('gray', '#9E9E9E'),
                            self.cores.get('grayDark', '#757575')
                        ]
            except Exception as e:
                print(f"Erro ao carregar cores do cliente: {e}")
        
        # Roda de Ouro (client_id == 4) - cores específicas (mais preto e cinza, menos dourado)
        if self.client_id == 4:
            self.cores.update({
                'primary': '#000000',  # Preto
                'primaryDark': '#000000',
                'primaryLight': '#1a1a1a',  # Preto mais claro
                'secondary': '#808080',  # Cinza (substitui dourado como secundária)
                'secondaryDark': '#606060',  # Cinza escuro
                'secondaryLight': '#A0A0A0',  # Cinza claro
                'masculino': '#000000',  # Preto
                'feminino': '#FFD700',  # Dourado (apenas para feminino)
                'preto': '#000000',
                'dourado': '#FFD700'  # Mantém dourado apenas para acentos
            })
            # Paleta: principalmente preto e cinzas, com dourado apenas como acento
            self.paleta = [
                '#000000',  # Preto
                '#1a1a1a',  # Preto claro
                '#404040',  # Cinza escuro
                '#606060',  # Cinza médio escuro
                '#808080',  # Cinza médio
                '#A0A0A0',  # Cinza claro
                '#C0C0C0',  # Cinza muito claro
                '#E0E0E0',  # Cinza quase branco
                '#FFD700',  # Dourado (apenas como acento)
                '#DAA520'   # Dourado escuro (apenas como acento)
            ]
    
    def _setup_custom_styles(self):
        """Configura estilos customizados - usa cores do cliente"""
        # Define cor primária (será atualizada após carregar cores)
        cor_primaria = self.cores.get('primary', '#1a237e') if hasattr(self, 'cores') else '#1a237e'
        cor_secundaria = self.cores.get('primaryLight', '#283593') if hasattr(self, 'cores') else '#283593'
        
        # Título principal
        self.title_style = ParagraphStyle(
            'CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor(cor_primaria),
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        # Subtítulo
        self.subtitle_style = ParagraphStyle(
            'CustomSubtitle',
            parent=self.styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor(cor_secundaria),
            spaceAfter=20,
            alignment=TA_CENTER
        )
        
        # Cabeçalho de seção
        self.section_style = ParagraphStyle(
            'CustomSection',
            parent=self.styles['Heading2'],
            fontSize=18,
            textColor=colors.HexColor(cor_primaria),
            spaceAfter=15,
            spaceBefore=20,
            fontName='Helvetica-Bold'
        )
        
        # Texto normal
        self.normal_style = ParagraphStyle(
            'CustomNormal',
            parent=self.styles['Normal'],
            fontSize=11,
            textColor=colors.HexColor('#212121'),
            spaceAfter=12
        )
        
        # Atualiza estilos após carregar cores (se já carregou)
        if hasattr(self, 'cores'):
            self._update_styles_colors()
    
    def _update_styles_colors(self):
        """Atualiza cores dos estilos após carregar cores do cliente"""
        cor_primaria = self.cores.get('primary', '#1a237e')
        cor_secundaria = self.cores.get('primaryLight', '#283593')
        
        self.title_style.textColor = colors.HexColor(cor_primaria)
        self.subtitle_style.textColor = colors.HexColor(cor_secundaria)
        self.section_style.textColor = colors.HexColor(cor_primaria)
    
    def _buscar_insight_grafico(self, tipo_grafico: str, insights: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
        """Busca o insight correspondente a um tipo de gráfico"""
        if not insights:
            return None
        
        # Mapeamento de tipos de gráfico para palavras-chave de busca nos títulos
        palavras_chave = {
            'evolucao_mensal': ['Evolução Mensal', 'evolucao', 'evolução'],
            'top_cids': ['TOP 10 Doenças', 'TOP 10 Doenças Mais Frequentes', 'CID', 'top_cids', 'doenças'],
            'funcionarios_dias': ['Dias Perdidos por Funcionário', 'funcionario', 'funcionário'],
            'top_setores': ['TOP 5 Setores', 'TOP 10 Setores', 'Setores', 'top_setores', 'setor'],
            'genero': ['Distribuição por Gênero', 'Gênero', 'genero', 'Masculino', 'Feminino'],
            'dias_doenca': ['Dias por Doença', 'Dias por', 'dias_doenca', 'doença'],
            'escalas': ['Escalas', 'escalas'],
            'motivos': ['Motivos', 'motivos', 'Incidência'],
            'centro_custo': ['Centro de Custo', 'centro_custo', 'Centro'],
            'distribuicao_dias': ['Distribuição de Dias', 'distribuicao_dias', 'Distribuição de Dias por Atestado'],
            'media_cid': ['Média de Dias por CID', 'media_cid', 'Média por CID'],
            'setor_genero': ['Setor e Gênero', 'setor_genero', 'Setor e']
        }
        
        if tipo_grafico not in palavras_chave:
            return None
        
        chaves_busca = palavras_chave[tipo_grafico]
        
        # Busca por título que contém qualquer uma das palavras-chave
        for insight in insights:
            titulo = insight.get('titulo', '').lower()
            # Remove prefixo "Análise: " se existir para facilitar busca
            titulo_limpo = titulo.replace('análise:', '').replace('analise:', '').strip()
            
            for chave in chaves_busca:
                if chave.lower() in titulo or chave.lower() in titulo_limpo:
                    return insight
        
        # Se não encontrou, retorna None
        return None
    
    def _gerar_grafico_cids(self, top_cids: List[Dict], output_path: str) -> Optional[str]:
        """Gera gráfico de barras horizontal para TOP CIDs"""
        try:
            if not top_cids or len(top_cids) == 0:
                return None
            
            fig, ax = plt.subplots(figsize=(10, 6))
            top10 = top_cids[:10]
            
            cids = [f"{c.get('cid', 'N/A')}" for c in top10]
            quantidades = [c.get('quantidade', 0) for c in top10]
            
            # Usa cores do cliente (preto e cinza para Roda de Ouro)
            colors_list = [self.cores['primary'] if i % 2 == 0 else self.cores['secondary'] for i in range(len(cids))]
            bars = ax.barh(range(len(cids)), quantidades, color=colors_list)
            
            ax.set_yticks(range(len(cids)))
            ax.set_yticklabels(cids, fontsize=9)
            ax.set_xlabel('Quantidade de Atestados', fontsize=10, fontweight='bold')
            ax.set_title('TOP 10 Doenças Mais Frequentes', fontsize=12, fontweight='bold', pad=15)
            ax.grid(axis='x', alpha=0.3)
            
            # Adiciona valores nas barras
            for i, (bar, qtd) in enumerate(zip(bars, quantidades)):
                ax.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2, 
                       str(qtd), va='center', fontsize=9)
            
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico CIDs: {e}")
            return None
    
    def _gerar_grafico_funcionarios(self, top_func: List[Dict], output_path: str) -> Optional[str]:
        """Gera gráfico de barras para TOP Funcionários"""
        try:
            if not top_func or len(top_func) == 0:
                return None
            
            fig, ax = plt.subplots(figsize=(10, 6))
            top10 = top_func[:10]
            
            nomes = [f.get('nome', 'N/A')[:20] for f in top10]
            dias = [f.get('dias_perdidos', 0) for f in top10]
            
            # Usa cores do cliente (para Roda de Ouro: preto e cinza)
            if self.client_id == 4:
                colors_list = [self.cores['primary'] if i % 2 == 0 else self.cores['secondary'] for i in range(len(nomes))]
            else:
                colors_list = [self.cores['primary'] if i % 2 == 0 else self.cores['secondary'] for i in range(len(nomes))]
            bars = ax.bar(range(len(nomes)), dias, color=colors_list)
            
            ax.set_xticks(range(len(nomes)))
            ax.set_xticklabels(nomes, rotation=45, ha='right', fontsize=9)
            ax.set_ylabel('Dias Perdidos', fontsize=10, fontweight='bold')
            ax.set_title('TOP 10 Funcionários com Mais Dias Perdidos', fontsize=12, fontweight='bold', pad=15)
            ax.grid(axis='y', alpha=0.3)
            
            # Adiciona valores nas barras
            for bar, dia in zip(bars, dias):
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height + 0.5,
                       str(int(dia)), ha='center', va='bottom', fontsize=9)
            
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico funcionários: {e}")
            return None
    
    def _gerar_grafico_setores(self, top_setores: List[Dict], output_path: str) -> Optional[str]:
        """Gera gráfico de barras para TOP Setores"""
        try:
            if not top_setores or len(top_setores) == 0:
                return None
            
            fig, ax = plt.subplots(figsize=(10, 6))
            top10 = top_setores[:10]
            
            setores = [s.get('setor', 'N/A')[:20] for s in top10]
            dias = [s.get('dias_perdidos', 0) for s in top10]
            
            # Usa cores do cliente
            colors_list = [self.cores['primary'] if i % 2 == 0 else self.cores['secondary'] for i in range(len(setores))]
            bars = ax.bar(range(len(setores)), dias, color=colors_list)
            
            ax.set_xticks(range(len(setores)))
            ax.set_xticklabels(setores, rotation=45, ha='right', fontsize=9)
            ax.set_ylabel('Dias Perdidos', fontsize=10, fontweight='bold')
            ax.set_title('TOP 10 Setores com Mais Dias Perdidos', fontsize=12, fontweight='bold', pad=15)
            ax.grid(axis='y', alpha=0.3)
            
            # Adiciona valores nas barras
            for bar, dia in zip(bars, dias):
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height + 0.5,
                       str(int(dia)), ha='center', va='bottom', fontsize=9)
            
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico setores: {e}")
            return None
    
    def _gerar_grafico_evolucao(self, evolucao: List[Dict], output_path: str) -> Optional[str]:
        """Gera gráfico de linha para evolução mensal"""
        try:
            if not evolucao or len(evolucao) == 0:
                return None
            
            fig, ax = plt.subplots(figsize=(12, 6))
            
            meses = [e.get('mes', 'N/A')[-5:] if len(e.get('mes', '')) >= 5 else e.get('mes', 'N/A') for e in evolucao]
            dias = [e.get('dias_perdidos', 0) for e in evolucao]
            quantidades = [e.get('quantidade', 0) for e in evolucao]
            
            # Eixo Y esquerdo para dias perdidos
            ax.plot(range(len(meses)), dias, marker='o', linewidth=2, markersize=6, 
                   color=self.cores['primary'], label='Dias Perdidos')
            ax.fill_between(range(len(meses)), dias, alpha=0.3, color=self.cores['primary'])
            ax.set_ylabel('Dias Perdidos', fontsize=10, fontweight='bold', color=self.cores['primary'])
            ax.tick_params(axis='y', labelcolor=self.cores['primary'])
            
            # Eixo Y direito para quantidade de atestados
            ax2 = ax.twinx()
            ax2.plot(range(len(meses)), quantidades, marker='s', linewidth=2, markersize=6, 
                    color=self.cores['secondary'], label='Quantidade de Atestados')
            ax2.fill_between(range(len(meses)), quantidades, alpha=0.3, color=self.cores['secondary'])
            ax2.set_ylabel('Quantidade de Atestados', fontsize=10, fontweight='bold', color=self.cores['secondary'])
            ax2.tick_params(axis='y', labelcolor=self.cores['secondary'])
            
            ax.set_xticks(range(len(meses)))
            ax.set_xticklabels(meses, rotation=45, ha='right', fontsize=9)
            ax.set_title('Evolução Mensal de Dias Perdidos e Atestados', fontsize=12, fontweight='bold', pad=15)
            ax.grid(True, alpha=0.3)
            
            # Legenda combinada
            lines1, labels1 = ax.get_legend_handles_labels()
            lines2, labels2 = ax2.get_legend_handles_labels()
            ax.legend(lines1 + lines2, labels1 + labels2, loc='upper left')
            
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico evolução: {e}")
            return None
    
    def _gerar_grafico_genero(self, distribuicao_genero: List[Dict], output_path: str) -> Optional[str]:
        """Gera gráfico de rosca (doughnut) para distribuição por gênero"""
        try:
            if not distribuicao_genero or len(distribuicao_genero) == 0:
                return None
            
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # Ordena os dados para garantir ordem correta: Masculino primeiro, depois Feminino
            dados_ordenados = []
            for item in distribuicao_genero:
                genero = item.get('genero', '')
                if genero == 'M':
                    dados_ordenados.append(('Masculino', item.get('quantidade', 0), 'M'))
                elif genero == 'F':
                    dados_ordenados.append(('Feminino', item.get('quantidade', 0), 'F'))
                else:
                    dados_ordenados.append((genero, item.get('quantidade', 0), genero))
            
            # Ordena: Masculino primeiro, depois Feminino
            dados_ordenados.sort(key=lambda x: (0 if x[2] == 'M' else 1 if x[2] == 'F' else 2))
            
            labels = [d[0] for d in dados_ordenados]
            valores = [d[1] for d in dados_ordenados]
            generos = [d[2] for d in dados_ordenados]
            
            # Mapeia cores corretamente: Masculino = preto, Feminino = dourado
            colors_list = []
            for gen in generos:
                if gen == 'M':
                    colors_list.append(self.cores['masculino'])  # Preto
                elif gen == 'F':
                    colors_list.append(self.cores['feminino'])  # Dourado
                else:
                    colors_list.append(self.cores.get('primary', '#1a237e'))
            
            wedges, texts, autotexts = ax.pie(valores, labels=labels, autopct='%1.1f%%',
                                             colors=colors_list, startangle=90,
                                             pctdistance=0.85)
            
            # Cria efeito de rosca
            centre_circle = plt.Circle((0, 0), 0.70, fc='white')
            ax.add_artist(centre_circle)
            
            # Adiciona total no centro
            total = sum(valores)
            ax.text(0, 0, f'Total\n{int(total)}', ha='center', va='center', 
                   fontsize=14, fontweight='bold', color=self.cores['primary'])
            
            ax.set_title('Distribuição por Gênero', fontsize=12, fontweight='bold', pad=15)
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico gênero: {e}")
            return None
    
    def _gerar_grafico_dias_doenca(self, top_cids_dias: List[Dict], output_path: str) -> Optional[str]:
        """Gera gráfico de barras para dias por doença"""
        try:
            if not top_cids_dias or len(top_cids_dias) == 0:
                return None
            
            fig, ax = plt.subplots(figsize=(10, 6))
            top5 = top_cids_dias[:5]
            
            descricoes = [d.get('descricao', d.get('cid', 'N/A'))[:30] for d in top5]
            dias = [d.get('dias_perdidos', 0) for d in top5]
            
            # Usa cores do cliente
            colors_list = [self.cores['secondary'] if i % 2 == 0 else self.cores['secondaryLight'] for i in range(len(descricoes))]
            bars = ax.bar(range(len(descricoes)), dias, color=colors_list)
            
            ax.set_xticks(range(len(descricoes)))
            ax.set_xticklabels(descricoes, rotation=45, ha='right', fontsize=9)
            ax.set_ylabel('Dias Perdidos', fontsize=10, fontweight='bold')
            ax.set_title('Dias Perdidos por Doença (TOP 5)', fontsize=12, fontweight='bold', pad=15)
            ax.grid(axis='y', alpha=0.3)
            
            # Adiciona valores nas barras
            for bar, dia in zip(bars, dias):
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height + 0.5,
                       str(int(dia)), ha='center', va='bottom', fontsize=9)
            
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico dias por doença: {e}")
            return None
    
    def _gerar_grafico_escalas(self, top_escalas: List[Dict], output_path: str) -> Optional[str]:
        """Gera gráfico de barras horizontal para escalas"""
        try:
            if not top_escalas or len(top_escalas) == 0:
                return None
            
            fig, ax = plt.subplots(figsize=(10, 6))
            top10 = top_escalas[:10]
            
            escalas = [e.get('escala', 'N/A')[:25] for e in top10]
            quantidades = [e.get('quantidade', 0) for e in top10]
            
            # Usa paleta do cliente
            bars = ax.barh(range(len(escalas)), quantidades, color=self.paleta[:len(escalas)])
            
            ax.set_yticks(range(len(escalas)))
            ax.set_yticklabels(escalas, fontsize=9)
            ax.set_xlabel('Quantidade de Atestados', fontsize=10, fontweight='bold')
            ax.set_title('Escalas com Mais Atestados (TOP 10)', fontsize=12, fontweight='bold', pad=15)
            ax.grid(axis='x', alpha=0.3)
            
            # Adiciona valores nas barras
            for i, (bar, qtd) in enumerate(zip(bars, quantidades)):
                ax.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2, 
                       str(qtd), va='center', fontsize=9)
            
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico escalas: {e}")
            return None
    
    def _gerar_grafico_motivos(self, top_motivos: List[Dict], output_path: str) -> Optional[str]:
        """Gera gráfico de pizza melhorado para motivos"""
        try:
            if not top_motivos or len(top_motivos) == 0:
                return None
            
            # Usa figura maior para melhor visualização
            fig, ax = plt.subplots(figsize=(12, 8))
            top10 = top_motivos[:10]
            
            motivos = [m.get('motivo', 'N/A') for m in top10]
            quantidades = [m.get('quantidade', 0) for m in top10]
            total = sum(quantidades)
            
            # Cores alternadas
            # Usa paleta do cliente
            colors_list = self.paleta[:len(motivos)] if len(motivos) <= len(self.paleta) else (self.paleta * ((len(motivos) // len(self.paleta)) + 1))[:len(motivos)]
            
            # Calcula percentuais
            percentuais = [(q / total * 100) if total > 0 else 0 for q in quantidades]
            
            # Cria labels mais legíveis (motivo + quantidade)
            labels = []
            for motivo, qtd, pct in zip(motivos, quantidades, percentuais):
                # Limita tamanho do motivo para não ficar muito longo
                motivo_curto = motivo[:20] + '...' if len(motivo) > 20 else motivo
                labels.append(f"{motivo_curto}\n({qtd} - {pct:.1f}%)")
            
            # Gera gráfico de pizza
            wedges, texts, autotexts = ax.pie(
                quantidades, 
                labels=None,  # Labels vão na legenda
                autopct=lambda pct: f'{pct:.1f}%' if pct > 5 else '',  # Mostra % apenas se > 5%
                colors=colors_list[:len(motivos)],
                startangle=90,
                textprops={'fontsize': 10, 'fontweight': 'bold', 'color': 'white'},
                pctdistance=0.75,
                labeldistance=1.05,
                explode=[0.05 if pct > 10 else 0 for pct in percentuais]  # Destaca fatias maiores
            )
            
            # Legenda melhorada à direita com informações completas
            legend_labels = []
            for motivo, qtd, pct in zip(motivos, quantidades, percentuais):
                # Formata motivo (trunca se muito longo)
                motivo_formatado = motivo if len(motivo) <= 35 else motivo[:32] + '...'
                legend_labels.append(f"{motivo_formatado}\n{qtd} atestados ({pct:.1f}%)")
            
            # Cria legenda mais espaçada
            legend = ax.legend(wedges, legend_labels, 
                             loc='center left', 
                             bbox_to_anchor=(1.05, 0.5),
                             fontsize=9, 
                             frameon=True, 
                             fancybox=True, 
                             shadow=True,
                             title='Motivos',
                             title_fontsize=10,
                             title_fontweight='bold')
            
            # Melhora o título da legenda
            legend.get_title().set_color('#1a237e')
            
            ax.set_title('Distribuição de Motivos de Atestados (TOP 10)', 
                        fontsize=14, fontweight='bold', pad=20, color='#1a237e')
            
            # Adiciona informações no rodapé
            ax.text(0, -1.15, f'Total: {int(total)} atestados', 
                   ha='center', va='center', 
                   fontsize=11, fontweight='bold', color='#556B2F',
                   bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.5))
            
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico motivos: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _gerar_grafico_distribuicao_dias(self, distribuicao_dias: List[Dict], output_path: str) -> Optional[str]:
        """Gera gráfico de barras para distribuição de dias por atestado"""
        try:
            if not distribuicao_dias or len(distribuicao_dias) == 0:
                return None
            
            fig, ax = plt.subplots(figsize=(12, 6))
            
            faixas = [str(d.get('faixa', d.get('dias', 'N/A'))) for d in distribuicao_dias]
            quantidades = [d.get('quantidade', 0) for d in distribuicao_dias]
            
            # Usa cores do cliente
            colors_list = [self.cores['primaryLight'] if i % 2 == 0 else self.cores['secondary'] for i in range(len(faixas))]
            bars = ax.bar(range(len(faixas)), quantidades, color=colors_list)
            
            ax.set_xticks(range(len(faixas)))
            ax.set_xticklabels(faixas, rotation=45, ha='right', fontsize=9)
            ax.set_ylabel('Quantidade de Atestados', fontsize=10, fontweight='bold')
            ax.set_xlabel('Dias por Atestado', fontsize=10, fontweight='bold')
            ax.set_title('Distribuição de Dias por Atestado', fontsize=12, fontweight='bold', pad=15)
            ax.grid(axis='y', alpha=0.3)
            
            # Adiciona valores nas barras
            for bar, qtd in zip(bars, quantidades):
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height + 0.5,
                       str(int(qtd)), ha='center', va='bottom', fontsize=9)
            
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico distribuição de dias: {e}")
            return None
    
    def _gerar_grafico_media_cid(self, media_cid: List[Dict], output_path: str) -> Optional[str]:
        """Gera gráfico de barras horizontal para média de dias por CID"""
        try:
            if not media_cid or len(media_cid) == 0:
                return None
            
            fig, ax = plt.subplots(figsize=(10, 6))
            top10 = media_cid[:10]
            
            labels = [f"{m.get('cid', 'N/A')} - {m.get('diagnostico', m.get('descricao', ''))[:20]}" 
                     for m in top10]
            medias = [m.get('media_dias', 0) for m in top10]
            
            # Usa cores do cliente
            colors_list = [self.cores['primary'] if i % 2 == 0 else self.cores['secondary'] for i in range(len(labels))]
            bars = ax.barh(range(len(labels)), medias, color=colors_list)
            
            ax.set_yticks(range(len(labels)))
            ax.set_yticklabels(labels, fontsize=9)
            ax.set_xlabel('Média de Dias', fontsize=10, fontweight='bold')
            ax.set_title('Média de Dias por CID (TOP 10)', fontsize=12, fontweight='bold', pad=15)
            ax.grid(axis='x', alpha=0.3)
            
            # Adiciona valores nas barras
            for i, (bar, media) in enumerate(zip(bars, medias)):
                ax.text(bar.get_width() + 0.1, bar.get_y() + bar.get_height()/2, 
                       f'{media:.1f}', va='center', fontsize=9)
            
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico média por CID: {e}")
            return None
    
    def _gerar_grafico_setor_genero(self, dias_setor_genero: List[Dict], output_path: str) -> Optional[str]:
        """Gera gráfico de barras agrupadas para setor e gênero"""
        try:
            if not dias_setor_genero or len(dias_setor_genero) == 0:
                return None
            
            # Agrupa dados por setor
            setores_map = {}
            for item in dias_setor_genero:
                setor = item.get('setor', 'N/A')
                if setor not in setores_map:
                    setores_map[setor] = {'M': 0, 'F': 0}
                genero = item.get('genero', '')
                dias = item.get('dias_perdidos', 0)
                if genero == 'M':
                    setores_map[setor]['M'] += dias
                elif genero == 'F':
                    setores_map[setor]['F'] += dias
            
            # Ordena por total de dias e pega top 10
            setores_ordenados = sorted(setores_map.items(), 
                                     key=lambda x: x[1]['M'] + x[1]['F'], 
                                     reverse=True)[:10]
            
            if not setores_ordenados:
                return None
            
            fig, ax = plt.subplots(figsize=(12, 6))
            
            setores = [s[0][:25] for s in setores_ordenados]
            masculino = [s[1]['M'] for s in setores_ordenados]
            feminino = [s[1]['F'] for s in setores_ordenados]
            
            x = np.arange(len(setores))
            width = 0.35
            
            # Usa cores do cliente (masculino e feminino)
            bars1 = ax.bar(x - width/2, masculino, width, label='Masculino', color=self.cores.get('masculino', '#1a237e'))
            bars2 = ax.bar(x + width/2, feminino, width, label='Feminino', color=self.cores.get('feminino', '#556B2F'))
            
            ax.set_xticks(x)
            ax.set_xticklabels(setores, rotation=45, ha='right', fontsize=9)
            ax.set_ylabel('Dias Perdidos', fontsize=10, fontweight='bold')
            ax.set_title('Dias Perdidos por Setor e Gênero (TOP 10)', fontsize=12, fontweight='bold', pad=15)
            ax.legend()
            ax.grid(axis='y', alpha=0.3)
            
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico setor e gênero: {e}")
            return None
    
    def _gerar_grafico_dias_ano_coerencia(self, dados: Dict, output_path: str) -> Optional[str]:
        """Gera gráfico de barras agrupadas para dias por ano e coerência"""
        try:
            if not dados or not isinstance(dados, dict):
                return None
            
            anos = dados.get('anos', [])
            coerente = dados.get('coerente', [])
            sem_coerencia = dados.get('sem_coerencia', [])
            
            if not anos or len(anos) == 0:
                return None
            
            fig, ax = plt.subplots(figsize=(12, 6))
            
            x = np.arange(len(anos))
            width = 0.35
            
            bars1 = ax.bar(x - width/2, coerente, width, label='Coerente', color=self.cores.get('primary', '#1a237e'))
            bars2 = ax.bar(x + width/2, sem_coerencia, width, label='Sem Coerência', color=self.cores.get('secondary', '#808080'))
            
            ax.set_xticks(x)
            ax.set_xticklabels(anos, fontsize=9)
            ax.set_ylabel('Dias Perdidos', fontsize=10, fontweight='bold')
            ax.set_title('Dias Atestados por Ano - Coerente vs Sem Coerência', fontsize=12, fontweight='bold', pad=15)
            ax.legend()
            ax.grid(axis='y', alpha=0.3)
            
            # Adiciona valores nas barras
            for bars in [bars1, bars2]:
                for bar in bars:
                    height = bar.get_height()
                    if height > 0:
                        ax.text(bar.get_x() + bar.get_width()/2., height,
                               f'{int(height)}', ha='center', va='bottom', fontsize=8)
            
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico dias por ano coerência: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _gerar_grafico_analise_coerencia(self, dados: Dict, output_path: str) -> Optional[str]:
        """Gera gráfico de rosca para análise de coerência"""
        try:
            if not dados or not isinstance(dados, dict):
                return None
            
            coerente = dados.get('coerente', 0)
            sem_coerencia = dados.get('sem_coerencia', 0)
            
            if coerente == 0 and sem_coerencia == 0:
                return None
            
            fig, ax = plt.subplots(figsize=(10, 6))
            
            labels = ['Coerente', 'Sem Coerência']
            valores = [coerente, sem_coerencia]
            colors_list = [self.cores.get('primary', '#1a237e'), self.cores.get('secondary', '#808080')]
            
            wedges, texts, autotexts = ax.pie(valores, labels=labels, autopct='%1.1f%%',
                                             colors=colors_list, startangle=90,
                                             pctdistance=0.85)
            
            # Cria efeito de rosca
            centre_circle = plt.Circle((0, 0), 0.70, fc='white')
            ax.add_artist(centre_circle)
            
            # Adiciona total no centro
            total = coerente + sem_coerencia
            ax.text(0, 0, f'Total\n{int(total)}', ha='center', va='center', 
                   fontsize=14, fontweight='bold', color=self.cores['primary'])
            
            ax.set_title('Análise Atestados - Coerente vs Sem Coerência', fontsize=12, fontweight='bold', pad=15)
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            return output_path
        except Exception as e:
            print(f"Erro ao gerar gráfico análise coerência: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def generate_pdf_report(self, 
                           output_path: str,
                           dados: Dict[str, Any],
                           metricas: Dict[str, Any],
                           insights: Optional[List[Dict[str, Any]]] = None,
                           periodo: Optional[str] = None,
                           insights_engine: Optional[Any] = None,
                           client_id: Optional[int] = None) -> bool:
        """Gera relatório PDF completo"""
        try:
            # Garante que o diretório existe
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            # Remove arquivo existente se houver (para evitar corrupção)
            if os.path.exists(output_path):
                try:
                    os.remove(output_path)
                except:
                    pass
            
            # SimpleDocTemplate será criado dentro do try/catch do build
            # (não criar aqui para evitar problemas de encoding)
            
            story = []
            
            # Função auxiliar para sanitizar texto (remove caracteres problemáticos)
            def sanitize_text(text):
                if not text:
                    return ""
                # Remove caracteres de controle que podem corromper PDF
                text = str(text)
                # Remove caracteres de controle (exceto \n, \r, \t)
                import re
                text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
                return text
            
            # Cabeçalho
            story.append(Paragraph(sanitize_text("Relatório de Absenteísmo"), self.title_style))
            story.append(Paragraph(sanitize_text("AbsenteismoController - GrupoBiomed"), self.subtitle_style))
            
            if periodo:
                story.append(Paragraph(sanitize_text(f"Período: {periodo}"), self.normal_style))
            
            story.append(Paragraph(sanitize_text(f"Data de geração: {datetime.now().strftime('%d/%m/%Y %H:%M')}"), self.normal_style))
            story.append(Spacer(1, 20))
            
            # Métricas principais
            story.append(Paragraph(sanitize_text("Indicadores Principais"), self.section_style))
            
            # Arredonda horas corretamente (sem truncar)
            total_horas = metricas.get('total_horas_perdidas', 0)
            horas_arredondadas = round(total_horas) if isinstance(total_horas, (int, float)) else int(total_horas or 0)
            
            metrics_data = [
                ['Métrica', 'Valor'],
                ['Total de Atestados', f"{int(metricas.get('total_atestados', 0))}"],
                ['Dias Perdidos', f"{int(metricas.get('total_dias_perdidos', 0))}"],
                ['Horas Perdidas', f"{horas_arredondadas}"],
                ['Atestados (Dias)', f"{int(metricas.get('total_atestados_dias', 0))}"],
            ]
            
            metrics_table = Table(metrics_data, colWidths=[10*cm, 6*cm])
            metrics_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.cores.get('primary', '#1a237e'))),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.grey),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
            ]))
            
            story.append(metrics_table)
            story.append(Spacer(1, 30))
            story.append(PageBreak())
            
            # Gera gráficos temporários
            temp_dir = os.path.dirname(output_path)
            timestamp_unique = datetime.now().strftime('%Y%m%d_%H%M%S_%f')
            graficos_temp = []
            
            # Define lista de gráficos a serem gerados (na ordem desejada)
            graficos_config = [
                ('evolucao_mensal', 'Evolução Mensal de Dias Perdidos e Atestados', self._gerar_grafico_evolucao),
                ('top_cids', 'TOP 10 Doenças Mais Frequentes', self._gerar_grafico_cids),
                ('top_funcionarios', 'TOP 10 Funcionários com Mais Dias Perdidos', self._gerar_grafico_funcionarios),
                ('top_setores', 'TOP 10 Setores com Mais Dias Perdidos', self._gerar_grafico_setores),
                ('distribuicao_genero', 'Distribuição por Gênero', self._gerar_grafico_genero),
                ('top_cids_dias', 'Dias Perdidos por Doença', self._gerar_grafico_dias_doenca),
                ('top_escalas', 'Escalas com Mais Atestados', self._gerar_grafico_escalas),
                ('top_motivos', 'Distribuição de Motivos de Atestados', self._gerar_grafico_motivos),
                ('dias_centro_custo', 'Dias Perdidos por Centro de Custo', self._gerar_grafico_setores),
                ('distribuicao_dias', 'Distribuição de Dias por Atestado', self._gerar_grafico_distribuicao_dias),
                ('media_cid', 'Média de Dias por CID', self._gerar_grafico_media_cid),
                ('dias_setor_genero', 'Dias Perdidos por Setor e Gênero', self._gerar_grafico_setor_genero),
            ]
            
            # Adiciona gráficos específicos da Roda de Ouro (client_id == 4) - REPLICA APRESENTAÇÃO
            if client_id == 4:
                graficos_config.extend([
                    ('classificacao_funcionarios_ro', 'Classificação por Funcionário', self._gerar_grafico_funcionarios),
                    ('classificacao_setores_ro', 'Classificação por Setor', self._gerar_grafico_setores),
                    ('classificacao_doencas_ro', 'Classificação por Doença', self._gerar_grafico_cids),
                    ('dias_ano_coerencia', 'Dias Atestados por Ano', self._gerar_grafico_dias_ano_coerencia),  # Função específica
                    ('analise_coerencia', 'Análise Atestados - Coerente vs Sem Coerência', self._gerar_grafico_analise_coerencia),  # Função específica
                    ('tempo_servico_atestados', 'Tempo Serviço x Atestados', self._gerar_grafico_setores),  # Replica apresentação
                    ('horas_perdidas_genero', 'Horas Perdidas por Gênero', self._gerar_grafico_genero),
                    ('horas_perdidas_setor', 'Horas Perdidas por Setor', self._gerar_grafico_setores),
                    ('evolucao_mensal_horas', 'Evolução Mensal de Horas Perdidas', self._gerar_grafico_evolucao),
                    ('comparativo_dias_horas_genero', 'Comparativo: Dias vs Horas vs Semanas', self._gerar_grafico_setor_genero),
                    ('horas_perdidas_setor_genero', 'Horas Perdidas por Setor e Gênero', self._gerar_grafico_setor_genero),
                    ('analise_detalhada_genero', 'Análise Detalhada por Gênero', self._gerar_grafico_genero),  # Replica apresentação
                ])
            
            # Gera cada gráfico com título e insight
            for chave_dados, titulo, funcao_grafico in graficos_config:
                if chave_dados in dados and dados[chave_dados]:
                    try:
                        # Valida se os dados não estão vazios ou em formato inválido
                        dados_grafico = dados[chave_dados]
                        if not dados_grafico:
                            print(f"⚠️ Dados vazios para {chave_dados}, pulando...")
                            continue
                        
                        # Valida formato dos dados (deve ser lista ou dict)
                        if not isinstance(dados_grafico, (list, dict)):
                            print(f"⚠️ Formato inválido para {chave_dados}: {type(dados_grafico)}, pulando...")
                            continue
                        
                        # Se for dict, verifica se não está vazio e se tem estrutura válida
                        if isinstance(dados_grafico, dict):
                            # Para dicts complexos da Roda de Ouro (ex: dias_ano_coerencia, analise_coerencia)
                            # Verifica se tem dados válidos
                            if len(dados_grafico) == 0:
                                print(f"⚠️ Dict vazio para {chave_dados}, pulando...")
                                continue
                            # Se for dict com estrutura específica (ex: {'anos': [], 'coerente': []})
                            # Verifica se tem pelo menos uma chave com dados
                            if any(isinstance(v, list) and len(v) > 0 for v in dados_grafico.values() if isinstance(v, list)):
                                # Tem dados, continua
                                pass
                            elif all(not v or (isinstance(v, (int, float)) and v == 0) for v in dados_grafico.values()):
                                print(f"⚠️ Dict sem dados válidos para {chave_dados}, pulando...")
                                continue
                        
                        # Se for list, verifica se não está vazio
                        if isinstance(dados_grafico, list) and len(dados_grafico) == 0:
                            print(f"⚠️ Lista vazia para {chave_dados}, pulando...")
                            continue
                        
                        grafico_path = os.path.join(temp_dir, f"{chave_dados}_{timestamp_unique}.png")
                        resultado_grafico = funcao_grafico(dados_grafico, grafico_path)
                        
                        if resultado_grafico and os.path.exists(grafico_path) and os.path.getsize(grafico_path) > 0:
                            graficos_temp.append(grafico_path)
                            
                            # Busca insight correspondente - tenta múltiplas estratégias
                            insight = None
                            if insights:
                                # Primeiro tenta buscar pelo mapeamento
                                insight = self._buscar_insight_grafico(chave_dados, insights)
                                
                                # Se não encontrou, tenta busca mais ampla por palavras-chave
                                if not insight:
                                    # Lista de palavras-chave alternativas para cada tipo
                                    busca_alternativa = {
                                        'evolucao_mensal': ['evolução', 'evolucao', 'mensal', 'tendência'],
                                        'top_cids': ['cid', 'doença', 'diagnóstico', 'frequente'],
                                        'funcionarios_dias': ['funcionário', 'funcionario', 'colaborador'],
                                        'top_setores': ['setor', 'departamento', 'área'],
                                        'genero': ['gênero', 'genero', 'masculino', 'feminino', 'sexo'],
                                        'dias_doenca': ['doença', 'dias perdidos', 'afastamento'],
                                        'escalas': ['escala', 'turno', 'horário'],
                                        'motivos': ['motivo', 'causa', 'razão', 'incidência'],
                                        'centro_custo': ['centro', 'custo', 'setor'],
                                        'distribuicao_dias': ['distribuição', 'distribuicao', 'histograma'],
                                        'media_cid': ['média', 'media', 'duração'],
                                        'setor_genero': ['setor', 'gênero', 'genero', 'comparativo']
                                    }
                                    
                                    if chave_dados in busca_alternativa:
                                        palavras = busca_alternativa[chave_dados]
                                        for ins in insights:
                                            titulo_lower = ins.get('titulo', '').lower()
                                            if any(palavra in titulo_lower for palavra in palavras):
                                                insight = ins
                                                break
                            
                            # Se ainda não encontrou insight e temos insights_engine, gera na hora
                            if not insight and insights_engine and chave_dados in dados and dados[chave_dados]:
                                try:
                                    analise_texto = insights_engine.gerar_analise_grafico(chave_dados, dados[chave_dados], metricas)
                                    if analise_texto:
                                        partes = analise_texto.split('💡')
                                        insight = {
                                            'tipo': 'analise',
                                            'icone': '📊',
                                            'titulo': f'Análise: {titulo}',
                                            'descricao': partes[0].strip().replace('**', '') if len(partes) > 0 else analise_texto.replace('**', ''),
                                            'recomendacao': partes[1].strip().replace('**', '').replace('💡', '').replace('Recomendação:', '').strip() if len(partes) > 1 else None
                                        }
                                except Exception as e:
                                    print(f"Erro ao gerar insight para {chave_dados}: {e}")
                            
                            # Cria conteúdo do gráfico (título + gráfico + insight) que deve ficar junto
                            conteudo_grafico = []
                            
                            # Título
                            conteudo_grafico.append(Paragraph(sanitize_text(titulo), self.section_style))
                            conteudo_grafico.append(Spacer(1, 10))
                            
                            # Gráfico (com verificação de existência e validação)
                            try:
                                if os.path.exists(grafico_path) and os.path.getsize(grafico_path) > 0:
                                    # Valida se é uma imagem válida (sem usar verify que fecha o arquivo)
                                    imagem_valida = True
                                    try:
                                        from PIL import Image as PILImage
                                        # Abre e fecha sem verify para não corromper o arquivo
                                        with PILImage.open(grafico_path) as pil_img:
                                            pil_img.load()  # Carrega a imagem sem fechar
                                        # Verifica se o arquivo é PNG válido lendo os primeiros bytes
                                        with open(grafico_path, 'rb') as f:
                                            header = f.read(8)
                                            if not (header.startswith(b'\x89PNG\r\n\x1a\n') or header.startswith(b'\xff\xd8\xff')):
                                                print(f"⚠️ Arquivo não é uma imagem PNG/JPG válida: {grafico_path}")
                                                imagem_valida = False
                                    except ImportError:
                                        # PIL não disponível, mas continua mesmo assim
                                        pass
                                    except Exception as img_error:
                                        print(f"⚠️ Erro ao validar imagem {grafico_path}: {img_error}")
                                        imagem_valida = False
                                    
                                    if imagem_valida:
                                        # Adiciona imagem ao PDF
                                        try:
                                            # Calcula proporção da imagem para manter aspect ratio
                                            try:
                                                from PIL import Image as PILImage
                                                with PILImage.open(grafico_path) as pil_img:
                                                    img_width, img_height = pil_img.size
                                                    aspect_ratio = img_width / img_height if img_height > 0 else 1.0
                                            except:
                                                aspect_ratio = 2.0  # Default 2:1
                                            
                                            # Ajusta altura baseado na largura e aspect ratio
                                            width_pdf = 16*cm
                                            height_pdf = width_pdf / aspect_ratio if aspect_ratio > 0 else 8*cm
                                            # Limita altura máxima
                                            if height_pdf > 10*cm:
                                                height_pdf = 10*cm
                                                width_pdf = height_pdf * aspect_ratio
                                            
                                            img = Image(grafico_path, width=width_pdf, height=height_pdf)
                                            conteudo_grafico.append(img)
                                            conteudo_grafico.append(Spacer(1, 15))
                                        except Exception as pdf_img_error:
                                            print(f"⚠️ Erro ao adicionar imagem ao PDF: {pdf_img_error}")
                                            import traceback
                                            traceback.print_exc()
                                            conteudo_grafico.append(Paragraph(f"<i>Gráfico não disponível</i>", self.normal_style))
                                            conteudo_grafico.append(Spacer(1, 15))
                                    else:
                                        conteudo_grafico.append(Paragraph(f"<i>Gráfico não disponível (imagem inválida)</i>", self.normal_style))
                                        conteudo_grafico.append(Spacer(1, 15))
                                else:
                                    conteudo_grafico.append(Paragraph(f"<i>Gráfico não disponível (arquivo não encontrado)</i>", self.normal_style))
                                    conteudo_grafico.append(Spacer(1, 15))
                            except Exception as e:
                                print(f"⚠️ Erro ao adicionar imagem {grafico_path} ao PDF: {e}")
                                import traceback
                                traceback.print_exc()
                                conteudo_grafico.append(Paragraph(f"<i>Erro ao carregar gráfico</i>", self.normal_style))
                                conteudo_grafico.append(Spacer(1, 15))
                            
                            # Insight abaixo do gráfico
                            if insight:
                                # Título do insight
                                insight_title_style = ParagraphStyle(
                                    'InsightTitle', 
                                    parent=self.normal_style, 
                                    fontSize=12, 
                                    textColor=colors.HexColor(self.cores.get('primary', '#1a237e')), 
                                    spaceAfter=8, 
                                    spaceBefore=10
                                )
                                # Remove emojis do título para evitar problemas no PDF
                                icone_texto = insight.get('icone', '📊')
                                # Tenta manter emoji, mas se der problema, remove
                                try:
                                    titulo_insight = f"<b>{icone_texto} Análise e Insights</b>"
                                except:
                                    titulo_insight = "<b>Análise e Insights</b>"
                                conteudo_grafico.append(Paragraph(
                                    sanitize_text(titulo_insight), 
                                    insight_title_style
                                ))
                        
                                # Descrição
                                descricao = sanitize_text(insight.get('descricao', '').replace('**', ''))
                                conteudo_grafico.append(Paragraph(descricao, self.normal_style))
                        
                                # Recomendação
                                if insight.get('recomendacao'):
                                    recomendacao = sanitize_text(insight.get('recomendacao').replace('**', ''))
                                    recomendacao_style = ParagraphStyle(
                                        'Recomendacao', 
                                        parent=self.normal_style, 
                                        fontSize=10, 
                                        leftIndent=20, 
                                        textColor=colors.HexColor('#556B2F'), 
                                        spaceAfter=15
                                    )
                                    # Remove emoji da recomendação
                                    conteudo_grafico.append(Paragraph(
                                        sanitize_text(f"<b>Recomendação:</b> {recomendacao}"), 
                                        recomendacao_style
                                    ))
                            
                            # Adiciona tudo junto usando KeepTogether para manter título e gráfico na mesma página
                            # Mas evita KeepTogether se o conteúdo for muito grande (pode causar problemas)
                            try:
                                # Conta quantos elementos há no conteúdo
                                num_elementos = len(conteudo_grafico)
                                # Se tiver muitos elementos, não usa KeepTogether para evitar problemas
                                if num_elementos <= 10:
                                    story.append(KeepTogether(conteudo_grafico))
                                else:
                                    # Adiciona sem KeepTogether se for muito grande
                                    for item in conteudo_grafico:
                                        story.append(item)
                            except Exception as e:
                                # Se KeepTogether falhar, adiciona sem ele
                                print(f"⚠️ Erro ao usar KeepTogether, adicionando conteúdo normalmente: {e}")
                                for item in conteudo_grafico:
                                    try:
                                        story.append(item)
                                    except Exception as item_error:
                                        print(f"⚠️ Erro ao adicionar item ao PDF: {item_error}")
                                        continue
                            story.append(Spacer(1, 20))
                    except Exception as e:
                        print(f"Erro ao processar gráfico {chave_dados}: {e}")
                        import traceback
                        traceback.print_exc()
                        continue
            
            story.append(PageBreak())
            
            story.append(Spacer(1, 20))
            
            # TOP 10 CIDs (Tabela)
            if 'top_cids' in dados and dados['top_cids']:
                story.append(Paragraph(sanitize_text("TOP 10 Doenças Mais Frequentes"), self.section_style))
                
                cids_data = [['CID', 'Diagnóstico', 'Quantidade', 'Dias Perdidos']]
                for cid in dados['top_cids'][:10]:
                    cids_data.append([
                        sanitize_text(cid.get('cid', 'N/A')),
                        sanitize_text(cid.get('descricao', cid.get('diagnostico', 'N/A'))[:50]),
                        str(cid.get('quantidade', 0)),
                        str(int(cid.get('dias_perdidos', 0)))
                    ])
                
                cids_table = Table(cids_data, colWidths=[3*cm, 7*cm, 3*cm, 3*cm])
                cids_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.cores.get('primary', '#1a237e'))),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('FONTSIZE', (0, 1), (-1, -1), 9),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.grey),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
                ]))
                
                story.append(cids_table)
                story.append(Spacer(1, 30))
            
            # TOP Funcionários
            if 'top_funcionarios' in dados and dados['top_funcionarios']:
                story.append(Paragraph(sanitize_text("TOP 10 Funcionários com Mais Dias Perdidos"), self.section_style))
                
                func_data = [['Funcionário', 'Setor', 'Dias Perdidos', 'Atestados']]
                for func in dados['top_funcionarios'][:10]:
                    func_data.append([
                        sanitize_text(func.get('nome', 'N/A')[:40]),
                        sanitize_text(func.get('setor', 'N/A')[:30]),
                        str(int(func.get('dias_perdidos', 0))),
                        str(func.get('quantidade_atestados', func.get('quantidade', 0)))
                    ])
                
                func_table = Table(func_data, colWidths=[5*cm, 4*cm, 3*cm, 3*cm])
                func_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.cores.get('primary', '#1a237e'))),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('FONTSIZE', (0, 1), (-1, -1), 9),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.grey),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
                ]))
                
                story.append(func_table)
                story.append(Spacer(1, 30))
            
            # TOP Setores
            if 'top_setores' in dados and dados['top_setores']:
                story.append(Paragraph(sanitize_text("TOP 10 Setores com Mais Atestados"), self.section_style))
                
                setor_data = [['Setor', 'Dias Perdidos', 'Atestados']]
                for setor in dados['top_setores'][:10]:
                    setor_data.append([
                        sanitize_text(setor.get('setor', 'N/A')[:40]),
                        str(int(setor.get('dias_perdidos', 0))),
                        str(setor.get('quantidade', 0))
                    ])
                
                setor_table = Table(setor_data, colWidths=[8*cm, 4*cm, 4*cm])
                setor_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.cores.get('primary', '#1a237e'))),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('FONTSIZE', (0, 1), (-1, -1), 9),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.grey),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
                ]))
                
                story.append(setor_table)
                story.append(Spacer(1, 30))
            
            # Rodapé
            story.append(Spacer(1, 20))
            story.append(Paragraph(
                sanitize_text("<i>Relatório gerado automaticamente pelo AbsenteismoController v2.0</i>"),
                self.normal_style
            ))
            
            # Valida se há conteúdo antes de gerar PDF
            if len(story) == 0:
                print("⚠️ Erro: Nenhum conteúdo para adicionar ao PDF!")
                return False
            
            # Build PDF com tratamento de erro robusto
            try:
                # Limpa qualquer arquivo existente antes de gerar
                if os.path.exists(output_path):
                    try:
                        os.remove(output_path)
                    except Exception as e:
                        print(f"⚠️ Aviso: Não foi possível remover arquivo existente: {e}")
                
                # Valida se há conteúdo antes de gerar
                if len(story) == 0:
                    print("❌ Erro: Nenhum conteúdo para adicionar ao PDF!")
                    return False
                
                # Gera PDF em arquivo temporário primeiro, depois move
                import tempfile
                import shutil
                temp_output = output_path + '.tmp'
                
                try:
                    # Cria SimpleDocTemplate (sem encoding, reportlab usa UTF-8 por padrão)
                    doc_temp = SimpleDocTemplate(
                        temp_output,
                        pagesize=A4,
                        rightMargin=2*cm,
                        leftMargin=2*cm,
                        topMargin=2*cm,
                        bottomMargin=2*cm
                    )
                    
                    # Gera PDF no arquivo temporário
                    doc_temp.build(story, onFirstPage=None, onLaterPages=None)
                    
                    # Fecha explicitamente
                    del doc_temp
                    
                    # Aguarda um pouco para garantir que o arquivo foi escrito
                    import time
                    time.sleep(0.3)
                    
                    # Valida arquivo temporário
                    if not os.path.exists(temp_output):
                        print(f"❌ Erro: PDF temporário não foi criado")
                        return False
                    
                    tamanho_temp = os.path.getsize(temp_output)
                    if tamanho_temp == 0:
                        print(f"❌ Erro: PDF temporário está vazio")
                        if os.path.exists(temp_output):
                            try:
                                os.remove(temp_output)
                            except:
                                pass
                        return False
                    
                    # Valida header do PDF temporário
                    with open(temp_output, 'rb') as f:
                        header = f.read(4)
                        if header != b'%PDF':
                            print(f"❌ Erro: Arquivo temporário não é PDF válido. Header: {header}")
                            if os.path.exists(temp_output):
                                try:
                                    os.remove(temp_output)
                                except:
                                    pass
                            return False
                    
                    # Move arquivo temporário para o destino final
                    shutil.move(temp_output, output_path)
                    
                    # Valida arquivo final
                    if not os.path.exists(output_path):
                        print(f"❌ Erro: PDF não foi movido para destino final")
                        return False
                    
                    tamanho = os.path.getsize(output_path)
                    if tamanho == 0:
                        print(f"❌ Erro: PDF final está vazio")
                        if os.path.exists(output_path):
                            try:
                                os.remove(output_path)
                            except:
                                pass
                        return False
                    
                    # Valida header do PDF final
                    with open(output_path, 'rb') as f:
                        header = f.read(8)
                        if not header.startswith(b'%PDF'):
                            print(f"❌ Erro: PDF final não tem header válido")
                            if os.path.exists(output_path):
                                try:
                                    os.remove(output_path)
                                except:
                                    pass
                            return False
                    
                    print(f"✅ PDF válido gerado: {output_path} ({tamanho} bytes)")
                    
                except Exception as build_inner_error:
                    print(f"❌ Erro interno ao construir PDF: {build_inner_error}")
                    import traceback
                    traceback.print_exc()
                    # Remove arquivos corrompidos
                    for file_path in [temp_output, output_path]:
                        if os.path.exists(file_path):
                            try:
                                os.remove(file_path)
                            except:
                                pass
                    return False
                
            except Exception as build_error:
                print(f"❌ Erro ao construir PDF: {build_error}")
                import traceback
                traceback.print_exc()
                # Remove arquivo corrompido se existir
                for file_path in [output_path, output_path + '.tmp']:
                    if os.path.exists(file_path):
                        try:
                            os.remove(file_path)
                        except:
                            pass
                return False
            
            # Remove arquivos temporários de gráficos
            for temp_file in graficos_temp:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except:
                    pass
            
            return True
            
        except Exception as e:
            print(f"Erro ao gerar PDF: {e}")
            import traceback
            traceback.print_exc()
            # Remove arquivo corrompido se existir
            if os.path.exists(output_path):
                try:
                    os.remove(output_path)
                except:
                    pass
            return False
    
    def generate_excel_report(self,
                             output_path: str,
                             dados: List[Dict[str, Any]],
                             metricas: Dict[str, Any],
                             dados_relatorio: Optional[Dict[str, Any]] = None,
                             periodo: Optional[str] = None,
                             client_id: Optional[int] = None) -> bool:
        """Gera relatório Excel completo com gráficos"""
        try:
            # Gera gráficos temporários primeiro
            temp_dir = os.path.dirname(output_path)
            timestamp_unique = datetime.now().strftime('%Y%m%d_%H%M%S_%f')
            graficos_temp = {}
            
            dados_graficos = dados_relatorio if dados_relatorio else {}
            
            # Define mapeamento de gráficos
            graficos_config = [
                ('evolucao_mensal', 'Evolução Mensal', self._gerar_grafico_evolucao),
                ('top_cids', 'TOP CIDs', self._gerar_grafico_cids),
                ('top_funcionarios', 'TOP Funcionários', self._gerar_grafico_funcionarios),
                ('top_setores', 'TOP Setores', self._gerar_grafico_setores),
                ('distribuicao_genero', 'Distribuição Gênero', self._gerar_grafico_genero),
                ('top_cids_dias', 'Dias por Doença', self._gerar_grafico_dias_doenca),
                ('top_escalas', 'Escalas', self._gerar_grafico_escalas),
                ('top_motivos', 'Motivos', self._gerar_grafico_motivos),
                ('dias_centro_custo', 'Centro de Custo', self._gerar_grafico_setores),
                ('distribuicao_dias', 'Distribuição Dias', self._gerar_grafico_distribuicao_dias),
                ('media_cid', 'Média por CID', self._gerar_grafico_media_cid),
                ('dias_setor_genero', 'Setor e Gênero', self._gerar_grafico_setor_genero),
            ]
            
            # Adiciona gráficos específicos da Roda de Ouro (client_id == 4) - REPLICA APRESENTAÇÃO
            if client_id == 4:
                graficos_config.extend([
                    ('classificacao_funcionarios_ro', 'Classificação Funcionários', self._gerar_grafico_funcionarios),
                    ('classificacao_setores_ro', 'Classificação Setores', self._gerar_grafico_setores),
                    ('classificacao_doencas_ro', 'Classificação Doenças', self._gerar_grafico_cids),
                    ('dias_ano_coerencia', 'Dias por Ano', self._gerar_grafico_evolucao),  # Replica apresentação
                    ('analise_coerencia', 'Análise Coerência', self._gerar_grafico_genero),  # Replica apresentação
                    ('tempo_servico_atestados', 'Tempo Serviço', self._gerar_grafico_setores),  # Replica apresentação
                    ('horas_perdidas_genero', 'Horas por Gênero', self._gerar_grafico_genero),
                    ('horas_perdidas_setor', 'Horas por Setor', self._gerar_grafico_setores),
                    ('evolucao_mensal_horas', 'Evolução Horas', self._gerar_grafico_evolucao),
                    ('comparativo_dias_horas_genero', 'Comparativo Dias/Horas', self._gerar_grafico_setor_genero),
                    ('horas_perdidas_setor_genero', 'Horas Setor/Gênero', self._gerar_grafico_setor_genero),
                    ('analise_detalhada_genero', 'Análise Detalhada Gênero', self._gerar_grafico_genero),  # Replica apresentação
                ])
            
            # Gera todos os gráficos
            for chave, nome_aba, funcao_grafico in graficos_config:
                if chave in dados_graficos and dados_graficos[chave]:
                    grafico_path = os.path.join(temp_dir, f"{chave}_{timestamp_unique}.png")
                    if funcao_grafico(dados_graficos[chave], grafico_path):
                        graficos_temp[chave] = grafico_path
            
            # Cria o Excel
            with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                # Aba 1: Dados Completos
                if dados:
                    df_dados = pd.DataFrame(dados)
                    df_dados.to_excel(writer, sheet_name='Dados Completos', index=False)
                
                # Aba 2: Métricas
                df_metricas = pd.DataFrame([
                    {'Métrica': 'Total de Atestados', 'Valor': metricas.get('total_atestados', 0)},
                    {'Métrica': 'Dias Perdidos', 'Valor': metricas.get('total_dias_perdidos', 0)},
                    {'Métrica': 'Horas Perdidas', 'Valor': metricas.get('total_horas_perdidas', 0)},
                    {'Métrica': 'Atestados (Dias)', 'Valor': metricas.get('total_atestados_dias', 0)},
                ])
                df_metricas.to_excel(writer, sheet_name='Métricas', index=False)
                
                # Abas com dados e gráficos
                mapeamento_abas = {
                    'top_cids': 'TOP CIDs',
                    'top_funcionarios': 'TOP Funcionários',
                    'top_setores': 'TOP Setores',
                    'evolucao_mensal': 'Evolução Mensal',
                    'distribuicao_genero': 'Distribuição Gênero',
                    'top_cids_dias': 'Dias por Doença',
                    'top_escalas': 'Escalas',
                    'top_motivos': 'Motivos',
                    'dias_centro_custo': 'Centro de Custo',
                    'distribuicao_dias': 'Distribuição Dias',
                    'media_cid': 'Média por CID',
                    'dias_setor_genero': 'Setor e Gênero',
                }
                
                # Adiciona abas específicas da Roda de Ouro - REPLICA APRESENTAÇÃO
                if client_id == 4:
                    mapeamento_abas.update({
                        'classificacao_funcionarios_ro': 'Classificação Funcionários',
                        'classificacao_setores_ro': 'Classificação Setores',
                        'classificacao_doencas_ro': 'Classificação Doenças',
                        'dias_ano_coerencia': 'Dias por Ano',  # Replica apresentação
                        'analise_coerencia': 'Análise Coerência',  # Replica apresentação
                        'tempo_servico_atestados': 'Tempo Serviço',  # Replica apresentação
                        'horas_perdidas_genero': 'Horas por Gênero',
                        'horas_perdidas_setor': 'Horas por Setor',
                        'evolucao_mensal_horas': 'Evolução Horas',
                        'comparativo_dias_horas_genero': 'Comparativo Dias/Horas',
                        'horas_perdidas_setor_genero': 'Horas Setor/Gênero',
                        'analise_detalhada_genero': 'Análise Detalhada Gênero',  # Replica apresentação
                    })
                
                for chave, nome_aba in mapeamento_abas.items():
                    if chave in dados_graficos and dados_graficos[chave]:
                        try:
                            df = pd.DataFrame(dados_graficos[chave])
                            df.to_excel(writer, sheet_name=nome_aba, index=False)
                        except Exception as e:
                            print(f"Erro ao criar aba {nome_aba}: {e}")
                            continue
            
            # Adiciona gráficos às abas e formatação
            try:
                from openpyxl import load_workbook
                from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
                from openpyxl.utils import get_column_letter
                
                wb = load_workbook(output_path)
                
                # Define cores do cliente
                cor_primaria_hex = self.cores.get('primary', '#1a237e')
                cor_secundaria_hex = self.cores.get('secondary', '#556B2F')
                
                # Converte hex para RGB
                def hex_to_rgb(hex_color):
                    hex_color = hex_color.lstrip('#')
                    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                
                cor_primaria_rgb = hex_to_rgb(cor_primaria_hex)
                cor_secundaria_rgb = hex_to_rgb(cor_secundaria_hex)
                
                # Estilos
                header_fill = PatternFill(start_color=''.join(f'{c:02X}' for c in cor_primaria_rgb), 
                                         end_color=''.join(f'{c:02X}' for c in cor_primaria_rgb), 
                                         fill_type='solid')
                header_font = Font(bold=True, color='FFFFFF', size=11)
                border = Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                )
                
                for chave, nome_aba in mapeamento_abas.items():
                    if nome_aba in wb.sheetnames:
                        try:
                            ws = wb[nome_aba]
                            
                            # Formata cabeçalho
                            if ws.max_row > 0:
                                for cell in ws[1]:
                                    cell.fill = header_fill
                                    cell.font = header_font
                                    cell.alignment = Alignment(horizontal='center', vertical='center')
                                    cell.border = border
                                
                                # Ajusta largura das colunas
                                for col in ws.columns:
                                    max_length = 0
                                    col_letter = get_column_letter(col[0].column)
                                    for cell in col:
                                        try:
                                            if len(str(cell.value)) > max_length:
                                                max_length = len(str(cell.value))
                                        except:
                                            pass
                                    adjusted_width = min(max_length + 2, 50)
                                    ws.column_dimensions[col_letter].width = adjusted_width
                                
                                # Aplica bordas nas células com dados
                                for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
                                    for cell in row:
                                        if cell.value:
                                            cell.border = border
                            
                            # Adiciona gráfico se disponível
                            if chave in graficos_temp:
                                try:
                                    img_path = graficos_temp[chave]
                                    if os.path.exists(img_path):
                                        img = ExcelImage(img_path)
                                        # Redimensiona para caber melhor no Excel
                                        img.width = 600
                                        img.height = int(img.height * (600 / img.width)) if img.width > 0 else 400
                                        
                                        # Posiciona após os dados (coluna A, após última linha)
                                        max_row = ws.max_row
                                        img_row = max(5, max_row + 3)
                                        img_cell = f'A{img_row}'
                                        ws.add_image(img, img_cell)
                                        
                                        # Ajusta altura da linha para acomodar a imagem
                                        ws.row_dimensions[img_row].height = img.height * 0.75
                                except Exception as e:
                                    print(f"Erro ao adicionar gráfico {nome_aba}: {e}")
                                    continue
                            
                        except Exception as e:
                            print(f"Erro ao formatar aba {nome_aba}: {e}")
                            continue
                
                wb.save(output_path)
                wb.close()
            except Exception as e:
                print(f"Erro ao adicionar gráficos ao Excel: {e}")
                import traceback
                traceback.print_exc()
            
            # Remove arquivos temporários
            for img_path in graficos_temp.values():
                try:
                    if os.path.exists(img_path):
                        os.remove(img_path)
                except:
                    pass
            
            return True
            
        except Exception as e:
            print(f"Erro ao gerar Excel: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def generate_powerpoint_report(self, output_path: str, dados_relatorio: Dict[str, Any], 
                                   metricas: Dict[str, Any], insights: List[Dict[str, Any]], 
                                   periodo: str = None, insights_engine=None, client_id: Optional[int] = None) -> bool:
        """Gera relatório em formato PowerPoint com gráficos e análises"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            
            # Cria apresentação
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Define cores da empresa (usa cores do cliente)
            cor_primaria_hex = self.cores.get('primary', '#1a237e')
            cor_secundaria_hex = self.cores.get('secondary', '#556B2F')
            
            # Para Roda de Ouro: usa cinza como secundária (não dourado)
            if client_id == 4:
                cor_secundaria_hex = '#808080'  # Cinza médio
            
            # Converte hex para RGB
            def hex_to_rgb(hex_color):
                hex_color = hex_color.lstrip('#')
                return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            
            cor_primaria = RGBColor(*hex_to_rgb(cor_primaria_hex))
            cor_secundaria = RGBColor(*hex_to_rgb(cor_secundaria_hex))
            cor_texto = RGBColor(33, 33, 33)  # #212121
            
            # Cor dourada apenas para acentos (Roda de Ouro)
            cor_dourado = RGBColor(255, 215, 0) if client_id == 4 else cor_secundaria
            
            # Cria diretório temporário para gráficos
            temp_dir = os.path.join(os.path.dirname(output_path), 'temp_pptx')
            os.makedirs(temp_dir, exist_ok=True)
            
            # Slide 0: Capa
            slide_capa = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco
            
            # Linha decorativa superior (usa cor secundária do cliente)
            linha_superior = slide_capa.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
            linha_superior.fill.solid()
            linha_superior.fill.fore_color.rgb = cor_secundaria
            linha_superior.line.color.rgb = cor_secundaria
            
            # Título principal
            caixa_titulo = slide_capa.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            frame_titulo = caixa_titulo.text_frame
            frame_titulo.text = "SAÚDE CORPORATIVA - MEDICINA DO TRABALHO"
            p_titulo = frame_titulo.paragraphs[0]
            p_titulo.font.size = Pt(32)
            p_titulo.font.bold = True
            p_titulo.font.color.rgb = cor_primaria
            p_titulo.alignment = PP_ALIGN.CENTER
            
            # Subtítulo
            caixa_subtitulo = slide_capa.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
            frame_subtitulo = caixa_subtitulo.text_frame
            frame_subtitulo.text = "INDICADORES DE SAÚDE"
            p_subtitulo = frame_subtitulo.paragraphs[0]
            p_subtitulo.font.size = Pt(28)
            p_subtitulo.font.bold = True
            p_subtitulo.font.color.rgb = cor_primaria
            p_subtitulo.alignment = PP_ALIGN.CENTER
            
            # Data
            data_str = periodo if periodo else datetime.now().strftime("%B, %Y")
            caixa_data = slide_capa.shapes.add_textbox(Inches(7), Inches(6.5), Inches(2.5), Inches(0.5))
            frame_data = caixa_data.text_frame
            frame_data.text = data_str.upper()
            p_data = frame_data.paragraphs[0]
            p_data.font.size = Pt(18)
            p_data.font.bold = True
            p_data.font.color.rgb = cor_primaria
            p_data.alignment = PP_ALIGN.RIGHT
            
            # Linha decorativa inferior (usa cor secundária do cliente)
            linha_inferior = slide_capa.shapes.add_shape(1, Inches(0), Inches(7.35), Inches(10), Inches(0.15))
            linha_inferior.fill.solid()
            linha_inferior.fill.fore_color.rgb = cor_secundaria
            linha_inferior.line.color.rgb = cor_secundaria
            
            dados_graficos = dados_relatorio if dados_relatorio else {}
            
            # Slide 1: KPIs
            if metricas:
                slide_kpis = prs.slides.add_slide(prs.slide_layouts[5])  # Layout Title Only
                title_kpis = slide_kpis.shapes.title
                title_kpis.text = "Indicadores Principais"
                title_kpis.text_frame.paragraphs[0].font.size = Pt(28)
                title_kpis.text_frame.paragraphs[0].font.color.rgb = cor_primaria
                title_kpis.text_frame.paragraphs[0].font.bold = True
                
                # Adiciona 3 KPIs
                kpi_values = [
                    ("Total de Atestados", int(metricas.get('total_atestados', 0) or metricas.get('total_registros', 0))),
                    ("Dias Perdidos", int(metricas.get('total_dias_perdidos', 0))),
                    ("Horas Perdidas", int(metricas.get('total_horas_perdidas', 0)))
                ]
                
                kpi_width = Inches(2.8)
                kpi_height = Inches(1.5)
                kpi_spacing = Inches(0.3)
                start_x = (Inches(10) - (kpi_width * 3 + kpi_spacing * 2)) / 2
                
                for i, (label, value) in enumerate(kpi_values):
                    left = start_x + i * (kpi_width + kpi_spacing)
                    top = Inches(2.5)
                    
                    # Retângulo de fundo
                    kpi_box = slide_kpis.shapes.add_shape(1, left, top, kpi_width, kpi_height)
                    kpi_box.fill.solid()
                    kpi_box.fill.fore_color.rgb = cor_primaria
                    kpi_box.line.color.rgb = cor_primaria
                    
                    # Valor
                    caixa_valor = slide_kpis.shapes.add_textbox(left + Inches(0.1), top + Inches(0.3), kpi_width - Inches(0.2), Inches(0.6))
                    frame_valor = caixa_valor.text_frame
                    frame_valor.text = str(value)
                    p_valor = frame_valor.paragraphs[0]
                    p_valor.font.size = Pt(36)
                    p_valor.font.bold = True
                    p_valor.font.color.rgb = RGBColor(255, 255, 255)
                    p_valor.alignment = PP_ALIGN.CENTER
                    
                    # Label
                    caixa_label = slide_kpis.shapes.add_textbox(left + Inches(0.1), top + Inches(1), kpi_width - Inches(0.2), Inches(0.4))
                    frame_label = caixa_label.text_frame
                    frame_label.text = label
                    p_label = frame_label.paragraphs[0]
                    p_label.font.size = Pt(12)
                    p_label.font.color.rgb = RGBColor(255, 255, 255)
                    p_label.alignment = PP_ALIGN.CENTER
                
                # Adiciona análise se disponível
                insight_kpis = None
                if insights:
                    for ins in insights:
                        if 'kpis' in ins.get('titulo', '').lower() or 'indicadores' in ins.get('titulo', '').lower():
                            insight_kpis = ins
                            break
                
                if not insight_kpis and insights_engine:
                    insight_kpis = {'texto': insights_engine.gerar_analise_grafico('kpis', None, metricas)}
                
                if insight_kpis:
                    texto_analise = insight_kpis.get('texto', '') if isinstance(insight_kpis, dict) else str(insight_kpis)
                    caixa_analise = slide_kpis.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
                    frame_analise = caixa_analise.text_frame
                    frame_analise.text = texto_analise[:500]  # Limita tamanho
                    frame_analise.word_wrap = True
                    p_analise = frame_analise.paragraphs[0]
                    p_analise.font.size = Pt(11)
                    p_analise.font.color.rgb = cor_texto
            
            # Mapeamento de gráficos para slides
            graficos_config = [
                ('evolucao_mensal', 'Evolução Mensal', self._gerar_grafico_evolucao),
                ('top_cids', 'TOP 10 Doenças mais Frequentes', self._gerar_grafico_cids),
                ('top_funcionarios', 'Dias Perdidos por Funcionário', self._gerar_grafico_funcionarios),
                ('top_setores', 'TOP 5 Setores', self._gerar_grafico_setores),
                ('distribuicao_genero', 'Distribuição por Gênero', self._gerar_grafico_genero),
                ('top_cids_dias', 'Dias por Doença', self._gerar_grafico_dias_doenca),
                ('top_escalas', 'Escalas com mais Atestados', self._gerar_grafico_escalas),
                ('top_motivos', 'Motivos de Incidência', self._gerar_grafico_motivos),
                ('dias_centro_custo', 'Dias Perdidos por Centro de Custo', self._gerar_grafico_setores),
                ('distribuicao_dias', 'Distribuição de Dias por Atestado', self._gerar_grafico_distribuicao_dias),
                ('media_cid', 'Média de Dias por CID', self._gerar_grafico_media_cid),
                ('dias_setor_genero', 'Dias Perdidos por Setor e Gênero', self._gerar_grafico_setor_genero),
            ]
            
            # Adiciona gráficos específicos da Roda de Ouro (client_id == 4) - REPLICA APRESENTAÇÃO
            if client_id == 4:
                graficos_config.extend([
                    ('classificacao_funcionarios_ro', 'Classificação por Funcionário', self._gerar_grafico_funcionarios),
                    ('classificacao_setores_ro', 'Classificação por Setor', self._gerar_grafico_setores),
                    ('classificacao_doencas_ro', 'Classificação por Doença', self._gerar_grafico_cids),
                    ('dias_ano_coerencia', 'Dias Atestados por Ano', self._gerar_grafico_dias_ano_coerencia),  # Função específica
                    ('analise_coerencia', 'Análise Atestados - Coerente vs Sem Coerência', self._gerar_grafico_analise_coerencia),  # Função específica
                    ('tempo_servico_atestados', 'Tempo Serviço x Atestados', self._gerar_grafico_setores),  # Replica apresentação
                    ('horas_perdidas_genero', 'Horas Perdidas por Gênero', self._gerar_grafico_genero),
                    ('horas_perdidas_setor', 'Horas Perdidas por Setor', self._gerar_grafico_setores),
                    ('evolucao_mensal_horas', 'Evolução Mensal de Horas Perdidas', self._gerar_grafico_evolucao),
                    ('comparativo_dias_horas_genero', 'Comparativo: Dias vs Horas vs Semanas', self._gerar_grafico_setor_genero),
                    ('horas_perdidas_setor_genero', 'Horas Perdidas por Setor e Gênero', self._gerar_grafico_setor_genero),
                    ('analise_detalhada_genero', 'Análise Detalhada por Gênero', self._gerar_grafico_genero),  # Replica apresentação
                ])
            
            # Gera slides para cada gráfico
            timestamp_unique = datetime.now().strftime('%Y%m%d_%H%M%S_%f')
            
            for chave, titulo, funcao_grafico in graficos_config:
                if chave in dados_graficos and dados_graficos[chave]:
                    # Gera gráfico
                    grafico_path = os.path.join(temp_dir, f"{chave}_{timestamp_unique}.png")
                    if funcao_grafico(dados_graficos[chave], grafico_path):
                        # Cria slide
                        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
                        
                        # Título
                        title_shape = slide.shapes.title
                        title_shape.text = titulo
                        title_shape.text_frame.paragraphs[0].font.size = Pt(24)
                        title_shape.text_frame.paragraphs[0].font.color.rgb = cor_primaria
                        title_shape.text_frame.paragraphs[0].font.bold = True
                        
                        # Adiciona gráfico
                        try:
                            slide.shapes.add_picture(grafico_path, Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(4))
                        except Exception as e:
                            print(f"Erro ao adicionar gráfico {chave}: {e}")
                        
                        # Busca insight
                        insight_grafico = self._buscar_insight_grafico(chave, insights)
                        if not insight_grafico and insights_engine:
                            insight_grafico = {'texto': insights_engine.gerar_analise_grafico(chave, dados_graficos[chave], metricas)}
                        
                        # Adiciona análise abaixo do gráfico
                        if insight_grafico:
                            texto_insight = insight_grafico.get('texto', '') if isinstance(insight_grafico, dict) else str(insight_grafico)
                            caixa_insight = slide.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(9), Inches(1.5))
                            frame_insight = caixa_insight.text_frame
                            frame_insight.text = texto_insight[:400]  # Limita tamanho
                            frame_insight.word_wrap = True
                            p_insight = frame_insight.paragraphs[0]
                            p_insight.font.size = Pt(10)
                            p_insight.font.color.rgb = cor_texto
                        
                        # Remove arquivo temporário
                        try:
                            if os.path.exists(grafico_path):
                                os.remove(grafico_path)
                        except:
                            pass
            
            # Salva apresentação
            prs.save(output_path)
            
            # Remove diretório temporário
            try:
                import shutil
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)
            except:
                pass
            
            return True
            
        except Exception as e:
            print(f"Erro ao gerar PowerPoint: {e}")
            import traceback
            traceback.print_exc()
            return False

